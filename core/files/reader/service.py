"""Оркестратор чтения файлов: единая точка входа FileReader.read -> FileReadResult."""

from __future__ import annotations

import asyncio
import base64
import mimetypes
import re
import shutil
import subprocess
import tempfile
import uuid
from collections import defaultdict
from collections.abc import Iterable, Mapping
from email import policy
from email.parser import BytesParser
from io import BytesIO
from pathlib import Path
from typing import Any, Protocol, TypeIs, cast

import ebooklib
import extract_msg
import fitz as pymupdf
import olefile
import trafilatura
import xlrd
from a2a.types import FilePart, FileWithBytes, Message, Part, Role, TextPart
from bs4 import BeautifulSoup
from charset_normalizer import from_bytes
from ebooklib import epub
from odf import opendocument as odf_opendocument
from odf import teletype as odf_teletype
from odf import text as odf_text
from pptx import Presentation as raw_presentation
from striprtf.striprtf import rtf_to_text
from unstructured.partition.auto import partition as raw_partition

from core.billing import get_billing_service
from core.billing.service import BALANCE_BLOCK_OPERATION_VISION
from core.clients.llm.factory import get_vision_llm
from core.context import get_context
from core.files.checksum import compute_content_checksum_sha256
from core.files.file_ref import FileRef, file_id_from_download_url, normalize_file_ref
from core.files.media.transcriber import MediaTranscriber
from core.files.models import FileRecord, FileResponse
from core.files.processors import get_default_file_processor
from core.files.reader.exceptions import FileReadError
from core.files.reader.models import (
    FileReadKind,
    FileReadResult,
    FileTypeInfo,
    ReadAsset,
    ReadAssetKind,
    ReadOptions,
    ReadPage,
    merge_file_ref_read_options,
)
from core.files.s3_client import S3ClientFactory
from core.files.types import FileCategory, extensions_for
from core.http import get_httpx_client
from core.models.billing_models import UsageType
from core.tracing.operation_span import traced_operation

SourceInput = Path | str | bytes


class _BytesReader(Protocol):
    def read(self, size: int = -1, /) -> bytes: ...


class _OleFileObject(Protocol):
    def exists(self, filename: str) -> bool: ...
    def openstream(self, filename: str) -> _BytesReader: ...
    def close(self) -> None: ...


class _OleFileModule(Protocol):
    def isOleFile(self, filename: object) -> bool: ...
    def OleFileIO(self, filename: object) -> _OleFileObject: ...


class _PdfPixmap(Protocol):
    width: int
    height: int

    def tobytes(self, output: str) -> bytes: ...


class _PdfPage(Protocol):
    def get_text(self) -> object: ...
    def get_pixmap(self, *, matrix: object) -> _PdfPixmap: ...


class _PdfDocument(Protocol):
    page_count: int

    def __getitem__(self, index: int) -> _PdfPage: ...
    def close(self) -> None: ...


class _FitzModule(Protocol):
    def open(self, *, stream: bytes, filetype: str) -> _PdfDocument: ...
    def Matrix(self, zoom_x: int, zoom_y: int) -> object: ...


class _UnstructuredMetadata(Protocol):
    page_number: object | None


class _UnstructuredElement(Protocol):
    metadata: _UnstructuredMetadata | None


class _UnstructuredPartition(Protocol):
    def __call__(
        self,
        *,
        file: object,
        metadata_filename: str,
        languages: list[str],
    ) -> Iterable[_UnstructuredElement]: ...


class _PptxRun(Protocol):
    text: str


class _PptxParagraph(Protocol):
    runs: Iterable[_PptxRun]


class _PptxTextFrame(Protocol):
    paragraphs: Iterable[_PptxParagraph]


class _PptxCell(Protocol):
    text: str


class _PptxRow(Protocol):
    cells: Iterable[_PptxCell]


class _PptxTable(Protocol):
    rows: Iterable[_PptxRow]


class _PptxShape(Protocol):
    has_text_frame: bool
    text_frame: _PptxTextFrame
    has_table: bool
    table: _PptxTable


class _PptxSlide(Protocol):
    shapes: Iterable[_PptxShape]


class _PptxPresentation(Protocol):
    slides: Iterable[_PptxSlide]


class _PptxPresentationFactory(Protocol):
    def __call__(self, file: object) -> _PptxPresentation: ...


class _OdfDocument(Protocol):
    def getElementsByType(self, elt: object) -> Iterable[object]: ...


class _OdfOpenDocumentModule(Protocol):
    def load(self, odffile: object) -> _OdfDocument: ...


class _OdfTextModule(Protocol):
    P: object
    H: object


class _OdfTeletypeModule(Protocol):
    def extractText(self, odfElement: object) -> str: ...


class _EpubItem(Protocol):
    def get_content(self) -> bytes: ...


class _EpubBook(Protocol):
    def get_items_of_type(self, item_type: object) -> Iterable[_EpubItem]: ...


class _EpubModule(Protocol):
    def read_epub(self, name: str, options: object | None = None) -> _EpubBook: ...


class _EbooklibModule(Protocol):
    ITEM_DOCUMENT: object


class _MsgFile(Protocol):
    subject: object | None
    sender: object | None
    to: object | None
    cc: object | None
    date: object | None
    body: object | None

    def close(self) -> None: ...


class _ExtractMsgModule(Protocol):
    def openMsg(self, path: str) -> _MsgFile: ...


def _is_file_ref_source(source: SourceInput | FileRef) -> TypeIs[FileRef]:
    if isinstance(source, (FileRecord, FileResponse)):
        return True
    if isinstance(source, Mapping) and not isinstance(source, (str, bytes, bytearray)):
        path_v = source.get("path")
        fid_v = source.get("file_id")
        url_v = source.get("url")
        if (
            (isinstance(path_v, str) and path_v.strip())
            or (isinstance(fid_v, str) and fid_v.strip())
            or (isinstance(url_v, str) and url_v.strip())
        ):
            return True
        raise TypeError(
            "Если source — словарь, укажите непустой path, file_id или url (запись вложения)."
        )
    return False


async def _read_local_file_bytes(path: Path) -> bytes:
    def _read() -> bytes:
        return path.read_bytes()

    return await asyncio.to_thread(_read)


async def _read_http_bytes(url: str) -> bytes:
    async with get_httpx_client(timeout=120.0) as client:
        response = await client.get(url)
    _ = response.raise_for_status()
    return response.content


def _s3_error_code(exc: Exception) -> str | None:
    response = cast(object, getattr(exc, "response", None))
    if not isinstance(response, Mapping):
        return None
    response_map = cast(Mapping[str, object], response)
    error = response_map.get("Error")
    if not isinstance(error, Mapping):
        return None
    error_map = cast(Mapping[str, object], error)
    code = error_map.get("Code")
    return code if isinstance(code, str) else None


async def _read_stored_file_by_id(file_id: str) -> tuple[bytes, str]:
    proc = await get_default_file_processor()
    record = await proc.get_file_record(file_id)
    if record is None:
        raise FileReadError(f"Файл не найден в хранилище: {file_id}")
    s3_bucket = getattr(record, "s3_bucket", None)
    s3_key = getattr(record, "s3_key", None)
    if isinstance(s3_bucket, str) and s3_bucket != "" and isinstance(s3_key, str) and s3_key != "":
        s3_client = S3ClientFactory.create_client_for_bucket(s3_bucket)
        try:
            try:
                raw = await s3_client.download_bytes(s3_key)
            except Exception as exc:
                code = _s3_error_code(exc)
                if code is None:
                    raise
                if code in ("NoSuchKey", "404", "NotFound"):
                    raise FileReadError(
                        "Файл не найден в хранилище: метаданные есть, объект отсутствует "
                        + "очистка бакета, смена окружения или устаревший идентификатор. "
                        + f"file_id={file_id}"
                    ) from exc
                raise FileReadError(
                    f"Ошибка объектного хранилища при чтении файла (код {code}): {file_id}"
                ) from exc
            return raw, record.original_name
        finally:
            await s3_client.close()
    storage_url = getattr(record, "storage_url", None)
    if isinstance(storage_url, str) and storage_url.startswith(("http://", "https://")):
        raw = await _read_http_bytes(storage_url)
        return raw, record.original_name
    raise FileReadError(f"Источник файла не настроен: {file_id}")


async def read_stored_file_by_id(file_id: str) -> tuple[bytes, str]:
    """Read a stored file by id and return raw bytes with the original file name."""
    return await _read_stored_file_by_id(file_id)


_TEXT_EXTENSIONS = extensions_for(FileCategory.TEXT)
_SPREADSHEET_EXTENSIONS = extensions_for(FileCategory.SPREADSHEET)
_OFFICE_EXTENSIONS = extensions_for(
    FileCategory.OFFICE_DOC,
    FileCategory.PRESENTATION,
    FileCategory.EMAIL,
    FileCategory.EBOOK,
)
_IMAGE_EXTENSIONS = extensions_for(FileCategory.IMAGE)
_AUDIO_EXTENSIONS = extensions_for(FileCategory.AUDIO)
_VIDEO_EXTENSIONS = extensions_for(FileCategory.VIDEO)

_DEFAULT_IMAGE_VISION_PROMPT = (
    "Извлеки весь видимый текст с изображения. "
    "Если текста нет, кратко опиши содержимое одним абзацем."
)
_TRANSCRIPTION_COMPANY_ID_REQUIRED = (
    "Транскрипция audio/video требует ReadOptions.transcription_company_id "
    + "или активной компании в контексте платформы."
)


def _normalize_extension(file_name: str) -> str:
    return Path(file_name).suffix.lower()


def _guess_mime(file_name: str) -> str | None:
    mime, _ = mimetypes.guess_type(file_name)
    return mime


def _payload_to_text(payload: object) -> str:
    if isinstance(payload, bytes):
        return payload.decode("utf-8", errors="replace")
    if isinstance(payload, str):
        return payload
    return ""


def _text_or_none(value: object | None) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    return text if text else None


def _mapping_text(mapping: Mapping[str, object], *keys: str) -> str | None:
    for key in keys:
        text = _text_or_none(mapping.get(key))
        if text is not None:
            return text
    return None


def _sniff_pdf(raw: bytes) -> bool:
    return len(raw) >= 5 and raw[:5] == b"%PDF-"


_OLE_COMPOUND_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _is_msword_ole_compound(raw: bytes) -> bool:
    """Бинарный Word 97–2003 (Compound File), ожидаемый antiword."""
    return (
        len(raw) >= len(_OLE_COMPOUND_MAGIC)
        and raw[: len(_OLE_COMPOUND_MAGIC)] == _OLE_COMPOUND_MAGIC
    )


def _is_zip_local_header_magic(raw: bytes) -> bool:
    """ZIP local file header (DOCX/XLSX и др. — OOXML под неверным расширением .doc)."""
    return len(raw) >= 4 and raw[:4] == b"PK\x03\x04"


def _try_antiword(raw: bytes) -> str | None:
    """Запускает antiword; возвращает текст или None при сбое."""
    antiword = shutil.which("antiword")
    if antiword is None:
        return None
    with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
        _ = tmp.write(raw)
        tmp_path = tmp.name
    try:
        proc = subprocess.run([antiword, tmp_path], capture_output=True, timeout=60)
    finally:
        Path(tmp_path).unlink(missing_ok=True)
    if proc.returncode != 0:
        return None
    text = proc.stdout.decode("utf-8", errors="replace").strip()
    return text or None


def _extract_text_from_ole_word_stream(raw: bytes) -> str | None:
    """
    Fallback извлечения текста из WordDocument stream OLE compound .doc.

    Не парсит структуру FIB полноценно (это сотни страниц спецификации MS-DOC),
    а вытаскивает читаемые ASCII/UTF-16 строки длиной >= 4 символов.
    Используется когда antiword не справляется (формат Word 95 / минимальные .doc).
    """
    typed_olefile = cast(_OleFileModule, cast(Any, olefile))

    if not typed_olefile.isOleFile(BytesIO(raw)):
        return None
    try:
        ole = typed_olefile.OleFileIO(BytesIO(raw))
    except Exception:
        return None
    try:
        if not ole.exists("WordDocument"):
            return None
        stream = ole.openstream("WordDocument").read()
    finally:
        ole.close()
    chunks: list[str] = []
    # ASCII-строки длиной >= 4 печатаемых символов
    for match in re.finditer(rb"[\x20-\x7e\r\n\t]{4,}", stream):
        chunks.append(match.group().decode("ascii", errors="replace"))
    # UTF-16 LE строки
    try:
        decoded16 = stream.decode("utf-16-le", errors="ignore")
        for match in re.finditer(r"[\u0020-\u00ff\u0400-\u04ff\u0100-\u017f]{4,}", decoded16):
            chunks.append(match.group())
    except Exception:
        pass
    text = "\n".join(chunks).strip()
    return text or None


def _kind_from_extension(ext: str) -> FileReadKind:
    if ext == ".pdf":
        return FileReadKind.PDF
    if ext in (".html", ".htm"):
        return FileReadKind.HTML
    if ext in _TEXT_EXTENSIONS:
        return FileReadKind.TEXT
    if ext in _SPREADSHEET_EXTENSIONS:
        return FileReadKind.SPREADSHEET
    if ext in _OFFICE_EXTENSIONS:
        return FileReadKind.OFFICE
    if ext in _IMAGE_EXTENSIONS:
        return FileReadKind.IMAGE
    if ext in _AUDIO_EXTENSIONS:
        return FileReadKind.AUDIO
    if ext in _VIDEO_EXTENSIONS:
        return FileReadKind.VIDEO
    return FileReadKind.UNKNOWN


class FileReader:
    """Чтение файлов в каноническую структуру FileReadResult."""

    def recognize_file_type(self, *, file_name: str, head: bytes | None = None) -> FileTypeInfo:
        ext = _normalize_extension(file_name)
        mime = _guess_mime(file_name)
        kind = _kind_from_extension(ext)
        if head is not None and _sniff_pdf(head):
            kind = FileReadKind.PDF
            if mime is None:
                mime = "application/pdf"
        if kind == FileReadKind.UNKNOWN and mime:
            if mime.startswith("image/"):
                kind = FileReadKind.IMAGE
            elif mime.startswith("audio/"):
                kind = FileReadKind.AUDIO
            elif mime.startswith("video/"):
                kind = FileReadKind.VIDEO
            elif mime == "application/pdf":
                kind = FileReadKind.PDF
            elif mime == "text/html":
                kind = FileReadKind.HTML
            elif mime in ("text/plain", "text/markdown", "text/csv"):
                kind = FileReadKind.TEXT
        return FileTypeInfo(detected_kind=kind, mime_type=mime, extension=ext)

    async def resolve_source(
        self,
        source: SourceInput,
        file_name: str | None,
        opts: ReadOptions,
    ) -> tuple[bytes, str]:
        if isinstance(source, bytes):
            if not file_name or not str(file_name).strip():
                raise ValueError("file_name обязателен при source=bytes")
            return source, str(file_name)

        path = source if isinstance(source, Path) else Path(str(source).strip())
        if path.is_file():
            data = await _read_local_file_bytes(path)
            name = file_name if file_name else path.name
            return data, name

        s = str(source).strip()
        if s.startswith(("http://", "https://")):
            data = await _read_http_bytes(s)
            tail = s.rsplit("/", 1)[-1]
            guessed = tail.split("?")[0] if tail else "file"
            name = file_name if file_name else (guessed or "file")
            return data, name

        fid: str | None = None
        if isinstance(opts.source_file_id, str) and opts.source_file_id.strip():
            fid = opts.source_file_id.strip()
        if not fid:
            fid = file_id_from_download_url(s)
        if fid:
            data, default_name = await _read_stored_file_by_id(fid)
            name = file_name if file_name else default_name
            return data, name

        raise FileReadError(f"Файл не найден: {s}")

    async def _read_resolved(
        self,
        raw: bytes,
        name: str,
        opts: ReadOptions,
    ) -> FileReadResult:
        source_checksum = opts.source_checksum or compute_content_checksum_sha256(raw)
        info = self.recognize_file_type(file_name=name, head=raw[:8192])
        mime = info.mime_type or _guess_mime(name)

        if info.detected_kind == FileReadKind.PDF or _sniff_pdf(raw):
            result = await asyncio.to_thread(_read_pdf_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.IMAGE:
            result = await _read_image_impl(raw, name, mime or "application/octet-stream", opts)
        elif info.detected_kind == FileReadKind.AUDIO:
            result = await _read_audio_impl(raw, name, mime or "audio/mpeg", opts)
        elif info.detected_kind == FileReadKind.VIDEO:
            result = await _read_video_impl(raw, name, mime or "video/mp4", opts)
        elif info.detected_kind == FileReadKind.HTML:
            result = await asyncio.to_thread(_read_html_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.TEXT:
            result = await asyncio.to_thread(_read_plain_text_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.SPREADSHEET and info.extension == ".xls":
            result = await asyncio.to_thread(_read_xls_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".doc":
            result = await asyncio.to_thread(_read_doc_choosing_backend_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".pptx":
            result = await asyncio.to_thread(_read_pptx_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".ppt":
            result = await asyncio.to_thread(_read_ppt_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".rtf":
            result = await asyncio.to_thread(_read_rtf_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".odt":
            result = await asyncio.to_thread(_read_odt_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".epub":
            result = await asyncio.to_thread(_read_epub_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".msg":
            result = await asyncio.to_thread(_read_msg_sync, raw, name, mime, opts)
        elif info.detected_kind == FileReadKind.OFFICE and info.extension == ".eml":
            result = await asyncio.to_thread(_read_eml_sync, raw, name, mime, opts)
        elif info.detected_kind in (FileReadKind.OFFICE, FileReadKind.SPREADSHEET):
            result = await asyncio.to_thread(
                _read_unstructured_sync, raw, name, mime, info.detected_kind, opts
            )
        elif info.detected_kind == FileReadKind.UNKNOWN:
            result = await asyncio.to_thread(
                _read_unstructured_sync, raw, name, mime, FileReadKind.UNKNOWN, opts
            )
        else:
            raise FileReadError(f"Неподдерживаемый тип: {info.detected_kind}")

        result = result.model_copy(
            update={
                "source_file_id": opts.source_file_id,
                "source_checksum": source_checksum,
            }
        )
        if result.page_count != len(result.pages):
            raise RuntimeError("инвариант FileReadResult: page_count должен совпадать с len(pages)")
        return result

    async def _raw_from_file_ref(
        self,
        finfo: Mapping[str, object],
        opts: ReadOptions,
    ) -> tuple[bytes, str]:
        display_name = _mapping_text(finfo, "name", "original_name")

        path_str = _mapping_text(finfo, "path")
        if path_str is not None:
            p = Path(path_str)
            if p.is_file():
                return await self.resolve_source(p, display_name, opts)

        url_val = _mapping_text(finfo, "url")
        if url_val is not None and url_val.startswith(("http://", "https://")):
            return await self.resolve_source(url_val, display_name, opts)

        source: SourceInput = ""
        if url_val is not None:
            source = url_val

        return await self.resolve_source(source, display_name, opts)

    async def read(
        self,
        source: SourceInput | FileRef,
        *,
        file_name: str | None = None,
        include_asset_bytes: bool = False,
        source_file_id: str | None = None,
        source_checksum: str | None = None,
        vision_model: str = "google/gemini-2.5-flash-preview",
        vision_prompt: str | None = None,
        transcription_company_id: str | None = None,
    ) -> FileReadResult:
        opts = ReadOptions(
            include_asset_bytes=include_asset_bytes,
            source_file_id=source_file_id,
            source_checksum=source_checksum,
            vision_model=vision_model,
            vision_prompt=vision_prompt,
            transcription_company_id=transcription_company_id,
        )
        if _is_file_ref_source(source):
            finfo = normalize_file_ref(source)
            opts = merge_file_ref_read_options(finfo, opts)
            raw, resolved_name = await self._raw_from_file_ref(finfo, opts)
            name = (
                file_name.strip() if isinstance(file_name, str) and file_name.strip() else None
            ) or resolved_name
        else:
            raw, name = await self.resolve_source(source, file_name, opts)
        return await self._read_resolved(raw, name, opts)


async def _read_image_impl(
    raw: bytes,
    file_name: str,
    mime: str,
    opts: ReadOptions,
) -> FileReadResult:
    b64 = base64.b64encode(raw).decode("utf-8")
    if opts.vision_prompt is not None:
        stripped = opts.vision_prompt.strip()
        if not stripped:
            raise ValueError("ReadOptions.vision_prompt задан пустой строкой")
        prompt = stripped
    else:
        prompt = _DEFAULT_IMAGE_VISION_PROMPT
    message = Message(
        message_id=str(uuid.uuid4()),
        role=Role.user,
        parts=[
            Part(root=TextPart(text=prompt)),
            Part(
                root=FilePart(
                    file=FileWithBytes(
                        bytes=b64,
                        mime_type=mime,
                        name=file_name,
                    )
                )
            ),
        ],
    )
    llm = get_vision_llm(model_name=opts.vision_model)
    actx = get_context()
    if actx is None or actx.active_company is None:
        raise ValueError("Контекст с active_company обязателен для vision-чтения изображения")
    if not str(actx.user.user_id).strip():
        raise ValueError(
            "Контекст с user обязателен для vision-чтения изображения (биллинг и уведомления)"
        )
    await get_billing_service().require_balance_for_billable_operation(
        actx.active_company.company_id,
        str(actx.user.user_id).strip(),
        operation_code=BALANCE_BLOCK_OPERATION_VISION,
        notification_service="frontend",
    )
    async with traced_operation(
        "core.files.reader.image",
        event_type="llm.vision",
        operation_category="llm",
        billing_usage_type=UsageType.LLM_REQUEST.value,
        billing_resource_name=f"llm:{opts.vision_model}",
        billing_quantity=1,
        billing_pending_settlement=True,
    ):
        vision_result = await llm.invoke([message], json_output=False)
    text = str(vision_result)
    page = ReadPage(index=0, text=text, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime,
        detected_kind=FileReadKind.IMAGE,
        page_count=1,
        pages=[page],
        warnings=[],
    )


def _read_html_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """HTML: trafilatura для контентных страниц + BeautifulSoup fallback для простой/частичной разметки."""
    del opts
    html = _decode_text_bytes(raw, file_name)
    try:
        extracted = trafilatura.extract(
            html,
            output_format="markdown",
            include_links=True,
            include_tables=True,
            include_images=False,
            favor_recall=True,
        )
    except Exception:
        extracted = None
    # Резерв на BeautifulSoup для простой/частичной HTML-разметки
    if extracted is None or extracted.strip() == "":
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        extracted = soup.get_text(separator="\n", strip=True)
    if not extracted or not extracted.strip():
        raise FileReadError(f"HTML не содержит текста: {file_name}")
    page = ReadPage(index=0, text=extracted, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "text/html",
        detected_kind=FileReadKind.HTML,
        page_count=1,
        pages=[page],
        warnings=[],
    )


def _decode_text_bytes(raw: bytes, file_name: str) -> str:
    """Декодирует текстовый файл с учётом BOM и авто-детекции кодировки."""
    # Порядок проверки BOM важен: 4-байтовые UTF-32 перед 2-байтовыми UTF-16
    boms = [
        (b"\xff\xfe\x00\x00", "utf-32-le"),
        (b"\x00\x00\xfe\xff", "utf-32-be"),
        (b"\xff\xfe", "utf-16-le"),
        (b"\xfe\xff", "utf-16-be"),
        (b"\xef\xbb\xbf", "utf-8-sig"),
    ]
    for bom_bytes, codec in boms:
        if raw.startswith(bom_bytes):
            return raw.decode(codec)
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        best = from_bytes(raw).best()
        if best is None:
            raise FileReadError(f"Не удалось определить кодировку текстового файла: {file_name}")
        return str(best)


def _read_plain_text_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    del opts
    text = _decode_text_bytes(raw, file_name)
    page = ReadPage(index=0, text=text, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "text/plain",
        detected_kind=FileReadKind.TEXT,
        page_count=1,
        pages=[page],
        warnings=[],
    )


def _read_xls_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    del opts
    book = xlrd.open_workbook(file_contents=raw)
    pages: list[ReadPage] = []
    for sheet_idx in range(book.nsheets):
        sheet = book.sheet_by_index(sheet_idx)
        rows: list[str] = []
        for row_idx in range(sheet.nrows):
            cells: list[str] = []
            for col_idx in range(sheet.ncols):
                cell = sheet.cell(row_idx, col_idx)
                value = cell.value
                if cell.ctype == xlrd.XL_CELL_EMPTY:
                    cells.append("")
                else:
                    cells.append(
                        str(value).rstrip("0").rstrip(".")
                        if isinstance(value, float) and value.is_integer()
                        else str(value)
                    )
            row_text = "\t".join(cells).rstrip()
            if row_text:
                rows.append(row_text)
        if rows:
            pages.append(
                ReadPage(index=sheet_idx, text="\n".join(rows), assets=[], label=sheet.name)
            )
    if not pages:
        raise FileReadError(f"xlrd не извлёк данные: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/vnd.ms-excel",
        detected_kind=FileReadKind.SPREADSHEET,
        page_count=len(pages),
        pages=pages,
        warnings=[],
    )


def _read_doc_choosing_backend_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """
    .doc:
      ZIP-сигнатура  → OOXML под расширением .doc → диспатч на Unstructured как .docx.
      OLE compound   → antiword; если он падает (часто на маленьких/legacy Word 95/97) → olefile-fallback.
      Иначе          → ошибка.
    """
    if _is_zip_local_header_magic(raw):
        eff_name = f"{file_name[:-4]}.docx" if file_name.lower().endswith(".doc") else file_name
        eff_mime = mime or "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        return _read_unstructured_sync(raw, eff_name, eff_mime, FileReadKind.OFFICE, opts)
    if _is_msword_ole_compound(raw):
        return _read_doc_ole_sync(raw, file_name, mime, opts)
    raise FileReadError(f"Файл .doc не распознан (нет OLE/ZIP сигнатур): {file_name}")


def _read_doc_ole_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """OLE compound .doc: antiword первым, затем olefile-fallback извлечения plain текста."""
    del opts
    text = _try_antiword(raw)
    if not text:
        text = _extract_text_from_ole_word_stream(raw)
    if not text or not text.strip():
        raise FileReadError(f".doc файл не содержит извлекаемого текста: {file_name}")
    page = ReadPage(index=0, text=text, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/msword",
        detected_kind=FileReadKind.OFFICE,
        page_count=1,
        pages=[page],
        warnings=[],
    )


def _read_pdf_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    fitz = cast(_FitzModule, cast(Any, pymupdf))

    warnings: list[str] = []
    doc = fitz.open(stream=raw, filetype="pdf")
    pages: list[ReadPage] = []
    encrypt_pdf_warning = False
    try:
        n = doc.page_count
        for i in range(n):
            skip_raster = False
            page = None
            try:
                page = doc[i]
                raw_text = page.get_text() or ""
                text = raw_text if isinstance(raw_text, str) else str(raw_text)
            except (ValueError, RuntimeError) as exc:
                msg = str(exc).lower()
                if "encrypt" in msg or "password" in msg:
                    if not encrypt_pdf_warning:
                        warnings.append("PDF зашифрован; текст недоступен без пароля")
                        encrypt_pdf_warning = True
                    text = ""
                    skip_raster = True
                else:
                    raise
            assets: list[ReadAsset] = []
            if opts.include_asset_bytes and not skip_raster:
                if page is None:
                    raise FileReadError(f"Не удалось получить страницу PDF: {file_name}#{i + 1}")
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                ch = compute_content_checksum_sha256(img_bytes)
                b64_val: str | None = None
                if opts.include_asset_bytes:
                    b64_val = base64.b64encode(img_bytes).decode("utf-8")
                assets.append(
                    ReadAsset(
                        kind=ReadAssetKind.PAGE_RASTER,
                        mime_type="image/png",
                        checksum=ch,
                        width=pix.width,
                        height=pix.height,
                        bytes_b64=b64_val,
                    )
                )
            pages.append(ReadPage(index=i, text=text, assets=assets, label=f"page_{i + 1}"))
    finally:
        doc.close()
    if not pages:
        warnings.append("PDF не содержит страниц")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/pdf",
        detected_kind=FileReadKind.PDF,
        page_count=len(pages),
        pages=pages,
        warnings=warnings,
    )


def _read_unstructured_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    kind: FileReadKind,
    opts: ReadOptions,
) -> FileReadResult:
    del opts
    partition = cast(_UnstructuredPartition, raw_partition)

    warnings: list[str] = []
    file_obj = BytesIO(raw)
    try:
        elements = partition(file=file_obj, metadata_filename=file_name, languages=["rus", "eng"])
    except Exception as exc:
        raise FileReadError(f"Не удалось разобрать файл через Unstructured: {file_name}") from exc
    by_page: dict[int, list[str]] = defaultdict(list)
    no_page_meta = False
    for el in elements:
        t = str(el).strip()
        if not t:
            continue
        md = el.metadata
        pn: int | None = None
        if md is not None:
            page_number = md.page_number
            if isinstance(page_number, int):
                pn = page_number
            elif isinstance(page_number, float) and page_number.is_integer():
                pn = int(page_number)
            elif isinstance(page_number, str) and page_number.strip():
                pn = int(page_number.strip())
        if pn is None:
            no_page_meta = True
            pn = 0
        else:
            pn = pn - 1
            if pn < 0:
                pn = 0
        by_page[pn].append(t)
    if not by_page:
        raise FileReadError(f"Unstructured не извлёк текст: {file_name}")
    if no_page_meta and len(by_page) == 1:
        warnings.append("Парсер не вернул номера страниц; весь текст на одной логической странице")
    sorted_keys = sorted(by_page.keys())
    pages: list[ReadPage] = []
    for seq, key in enumerate(sorted_keys):
        body = "\n\n".join(by_page[key])
        label = None if len(sorted_keys) == 1 else f"page_{key + 1}"
        pages.append(ReadPage(index=seq, text=body, assets=[], label=label))
    detected = kind
    if kind == FileReadKind.UNKNOWN:
        detected = FileReadKind.OFFICE
    return FileReadResult(
        file_name=file_name,
        mime_type=mime,
        detected_kind=detected,
        page_count=len(pages),
        pages=pages,
        warnings=warnings,
    )


def _read_pptx_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """PowerPoint .pptx через python-pptx; одна страница на слайд."""
    del opts
    presentation = cast(_PptxPresentationFactory, raw_presentation)
    try:
        prs = presentation(BytesIO(raw))
    except Exception as exc:
        raise FileReadError(f"Не удалось открыть PPTX: {file_name}") from exc
    pages: list[ReadPage] = []
    for slide_idx, slide in enumerate(prs.slides):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            chunks.append(run.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    chunks.append("\t".join(cells))
        body = "\n".join(c for c in chunks if c)
        pages.append(
            ReadPage(index=slide_idx, text=body, assets=[], label=f"slide_{slide_idx + 1}")
        )
    if not pages:
        raise FileReadError(f"PPTX не содержит слайдов: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime
        or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        detected_kind=FileReadKind.OFFICE,
        page_count=len(pages),
        pages=pages,
        warnings=[],
    )


def _read_ppt_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """PowerPoint legacy .ppt: ZIP-сигнатура → переименовать как .pptx, иначе сообщить."""
    del opts
    if _is_zip_local_header_magic(raw):
        return _read_pptx_sync(
            raw,
            file_name[:-4] + ".pptx" if file_name.lower().endswith(".ppt") else file_name,
            mime,
            ReadOptions(),
        )
    raise FileReadError(
        f"Legacy .ppt (PowerPoint 97-2003) не поддерживается без libreoffice: {file_name}. "
        + "Конвертируйте в .pptx."
    )


def _read_rtf_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """RTF через striprtf (чисто Python, без зависимостей)."""
    del opts
    try:
        rtf_text = raw.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_text, errors="ignore")
    except Exception as exc:
        raise FileReadError(f"Не удалось разобрать RTF: {file_name}") from exc
    if not text.strip():
        raise FileReadError(f"RTF пустой: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/rtf",
        detected_kind=FileReadKind.OFFICE,
        page_count=1,
        pages=[ReadPage(index=0, text=text, assets=[], label=None)],
        warnings=[],
    )


def _read_odt_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """OpenDocument Text .odt через odfpy (ZIP+XML)."""
    del opts
    teletype = cast(_OdfTeletypeModule, cast(Any, odf_teletype))
    odf_text_module = cast(_OdfTextModule, cast(Any, odf_text))
    opendocument = cast(
        _OdfOpenDocumentModule,
        cast(Any, odf_opendocument),
    )

    try:
        doc = opendocument.load(BytesIO(raw))
    except Exception as exc:
        raise FileReadError(f"Не удалось открыть ODT: {file_name}") from exc
    parts: list[str] = []
    for elem in doc.getElementsByType(odf_text_module.P):
        s = teletype.extractText(elem)
        if s.strip():
            parts.append(s)
    for elem in doc.getElementsByType(odf_text_module.H):
        s = teletype.extractText(elem)
        if s.strip():
            parts.append(s)
    body = "\n".join(parts)
    if not body.strip():
        raise FileReadError(f"ODT пустой: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/vnd.oasis.opendocument.text",
        detected_kind=FileReadKind.OFFICE,
        page_count=1,
        pages=[ReadPage(index=0, text=body, assets=[], label=None)],
        warnings=[],
    )


def _read_epub_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """EPUB через EbookLib + BeautifulSoup; одна страница на главу."""
    del opts
    typed_ebooklib = cast(_EbooklibModule, cast(Any, ebooklib))
    typed_epub = cast(_EpubModule, cast(Any, epub))

    # ebooklib читает только с диска, поэтому пишем во временный файл
    with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tmp:
        _ = tmp.write(raw)
        tmp_path = tmp.name
    try:
        try:
            book = typed_epub.read_epub(tmp_path)
        except Exception as exc:
            raise FileReadError(f"Не удалось открыть EPUB: {file_name}") from exc
        pages: list[ReadPage] = []
        for idx, item in enumerate(book.get_items_of_type(typed_ebooklib.ITEM_DOCUMENT)):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text:
                pages.append(
                    ReadPage(
                        index=len(pages),
                        text=text,
                        assets=[],
                        label=f"chapter_{idx + 1}",
                    )
                )
    finally:
        Path(tmp_path).unlink(missing_ok=True)
    if not pages:
        raise FileReadError(f"EPUB не содержит документов: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/epub+zip",
        detected_kind=FileReadKind.OFFICE,
        page_count=len(pages),
        pages=pages,
        warnings=[],
    )


def _read_msg_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """Outlook .msg через extract-msg (чисто Python)."""
    del opts
    typed_extract_msg = cast(_ExtractMsgModule, cast(Any, extract_msg))

    with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
        _ = tmp.write(raw)
        tmp_path = tmp.name
    try:
        try:
            msg = typed_extract_msg.openMsg(tmp_path)
        except Exception as exc:
            raise FileReadError(f"Не удалось открыть MSG: {file_name}") from exc
        try:
            parts: list[str] = []
            subject = _text_or_none(msg.subject)
            sender = _text_or_none(msg.sender)
            to = _text_or_none(msg.to)
            cc = _text_or_none(msg.cc)
            date = _text_or_none(msg.date)
            body = _text_or_none(msg.body)
            if subject:
                parts.append(f"Subject: {subject}")
            if sender:
                parts.append(f"From: {sender}")
            if to:
                parts.append(f"To: {to}")
            if cc:
                parts.append(f"Cc: {cc}")
            if date:
                parts.append(f"Date: {date}")
            if body:
                parts.append("")
                parts.append(body)
            text = "\n".join(parts)
        finally:
            msg.close()
    finally:
        Path(tmp_path).unlink(missing_ok=True)
    if not text.strip():
        raise FileReadError(f"MSG пустой: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "application/vnd.ms-outlook",
        detected_kind=FileReadKind.OFFICE,
        page_count=1,
        pages=[ReadPage(index=0, text=text, assets=[], label=None)],
        warnings=[],
    )


def _read_eml_sync(
    raw: bytes,
    file_name: str,
    mime: str | None,
    opts: ReadOptions,
) -> FileReadResult:
    """RFC 822 .eml через встроенный stdlib email модуль."""
    del opts
    try:
        msg = BytesParser(policy=policy.default).parsebytes(raw)
    except Exception as exc:
        raise FileReadError(f"Не удалось разобрать EML: {file_name}") from exc
    parts: list[str] = []
    for header in ("Subject", "From", "To", "Cc", "Date"):
        v = msg.get(header)
        if v:
            parts.append(f"{header}: {v}")
    parts.append("")
    body: str | None = None
    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            if ctype == "text/plain" and part.get("Content-Disposition") is None:
                try:
                    body = _payload_to_text(cast(object, part.get_content()))
                except (LookupError, UnicodeDecodeError):
                    body = _payload_to_text(part.get_payload(decode=True))
                break
        if body is None:
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/html":
                    try:
                        html = (
                            _payload_to_text(cast(object, part.get_content()))
                            if hasattr(part, "get_content")
                            else _payload_to_text(part.get_payload(decode=True))
                        )
                        body = BeautifulSoup(html, "html.parser").get_text(
                            separator="\n", strip=True
                        )
                    except Exception:
                        continue
                    break
    else:
        try:
            body = _payload_to_text(cast(object, msg.get_content()))
        except (LookupError, UnicodeDecodeError):
            body = _payload_to_text(msg.get_payload(decode=True))
    if body:
        parts.append(body)
    text = "\n".join(parts)
    if not text.strip():
        raise FileReadError(f"EML пустой: {file_name}")
    return FileReadResult(
        file_name=file_name,
        mime_type=mime or "message/rfc822",
        detected_kind=FileReadKind.OFFICE,
        page_count=1,
        pages=[ReadPage(index=0, text=text, assets=[], label=None)],
        warnings=[],
    )


def _resolve_transcription_company_id(opts: ReadOptions) -> str:
    if opts.transcription_company_id is not None and opts.transcription_company_id.strip() != "":
        return opts.transcription_company_id.strip()
    ctx = get_context()
    if ctx is None or ctx.active_company is None:
        raise ValueError(_TRANSCRIPTION_COMPANY_ID_REQUIRED)
    company_id = ctx.active_company.company_id
    if company_id == "":
        raise ValueError(_TRANSCRIPTION_COMPANY_ID_REQUIRED)
    return company_id


async def _read_audio_impl(
    raw: bytes,
    file_name: str,
    mime: str,
    opts: ReadOptions,
) -> FileReadResult:
    company_id = _resolve_transcription_company_id(opts)
    transcriber = MediaTranscriber(company_id=company_id)
    transcription = await transcriber.transcribe_audio(
        audio_bytes=raw,
        file_name=file_name,
        mime_type=mime,
    )
    page = ReadPage(index=0, text=transcription.text, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime,
        detected_kind=FileReadKind.AUDIO,
        page_count=1,
        pages=[page],
        warnings=[],
    )


async def _read_video_impl(
    raw: bytes,
    file_name: str,
    mime: str,
    opts: ReadOptions,
) -> FileReadResult:
    company_id = _resolve_transcription_company_id(opts)
    transcriber = MediaTranscriber(company_id=company_id)
    transcription = await transcriber.transcribe_video(
        video_bytes=raw,
        file_name=file_name,
    )
    page = ReadPage(index=0, text=transcription.text, assets=[], label=None)
    return FileReadResult(
        file_name=file_name,
        mime_type=mime,
        detected_kind=FileReadKind.VIDEO,
        page_count=1,
        pages=[page],
        warnings=[],
    )
